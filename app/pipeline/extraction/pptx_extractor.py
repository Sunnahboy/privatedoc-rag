from pathlib import Path

import anyio
from pptx import Presentation

from .base import BaseExtractor, ExtractionResult


class PPTExtractor(BaseExtractor):
    """
    Extract text from PowerPoint presentation
    """

    def _sync_extract(self, file_path: Path) -> ExtractionResult:
        presentation = Presentation(file_path)

        slides = []

        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "txt") and shape.text.strip():
                    slide.append(shape.txt)

        return ExtractionResult(
            text="\n".join(slides),
            page_count=len(presentation.slides),
            metadata={
                "slides": len(presentation.slides),
            },
        )

    async def extract(self, file_path) -> ExtractionResult:
        return await anyio.to_thread.run_sync(
            self._sync_extract,
            file_path,
        )
